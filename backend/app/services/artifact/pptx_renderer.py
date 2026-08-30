import os
import pptx
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

from app.schemas.artifact import ArtifactPlan, SlideModel, SlideType

class PPTXRenderer:
    def __init__(self):
        pass

    def render(self, plan: ArtifactPlan, output_path: str) -> str:
        """
        Renders the given ArtifactPlan into a .pptx file at the output_path.
        Returns the absolute output_path on success.
        Throws exceptions on failure and ensures no partial files are left.
        """
        tmp_path = f"{output_path}.tmp"
        try:
            prs = pptx.Presentation()

            for slide_model in plan.slides:
                self._render_slide(prs, slide_model)

            prs.save(tmp_path)
            
            # Atomic promotion
            if os.path.exists(output_path):
                os.remove(output_path)
            os.replace(tmp_path, output_path)
            
            return output_path
            
        except Exception as e:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
            raise e

    def _render_slide(self, prs: pptx.Presentation, slide_model: SlideModel):
        if slide_model.slide_type == SlideType.TITLE:
            layout_index = 0
        else:
            layout_index = 1
            
        # Get the layout
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title_shape = slide.shapes.title
        if title_shape and title_shape.text_frame:
            title_shape.text = slide_model.title
            title_shape.text_frame.word_wrap = True
            
        # Set Content (only for non-title layout)
        if layout_index == 1 and len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            if body_shape.text_frame:
                tf = body_shape.text_frame
                tf.word_wrap = True
                
                # Check for overloaded content
                total_chars = sum(len(c) for c in slide_model.content)
                if total_chars > 2000:
                    # Truncate slightly or just log, but the requirements say 
                    # "add reasonable deterministic text-density/length safeguards"
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                tf.text = "" # Clear default text
                for idx, content_bullet in enumerate(slide_model.content):
                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = content_bullet
                    p.level = 0
                    
        # Add Speaker Notes / Provenance
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        
        parts = []
        if slide_model.speaker_notes:
            parts.append(f"[Speaker Notes]\n{slide_model.speaker_notes}")
            
        if slide_model.source_node_ids or slide_model.evidence_ids:
            parts.append("[Provenance]")
            if slide_model.source_node_ids:
                parts.append(f"Sources: {', '.join(slide_model.source_node_ids)}")
            if slide_model.evidence_ids:
                parts.append(f"Evidence: {', '.join(slide_model.evidence_ids)}")
                
        if parts:
            notes_tf.text = "\n\n".join(parts)
