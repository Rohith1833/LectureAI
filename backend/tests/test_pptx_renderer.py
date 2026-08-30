import os
import tempfile
import pytest
import pptx
from app.schemas.artifact import ArtifactPlan, SlideModel, SlideType
from app.services.artifact.pptx_renderer import PPTXRenderer

class TestPPTXRenderer:
    @pytest.fixture
    def renderer(self):
        return PPTXRenderer()

    @pytest.fixture
    def temp_file(self):
        fd, path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        yield path
        if os.path.exists(path):
            os.remove(path)
        if os.path.exists(f"{path}.tmp"):
            os.remove(f"{path}.tmp")

    def test_render_minimal_plan(self, renderer, temp_file):
        plan = ArtifactPlan(slides=[
            SlideModel(slide_type=SlideType.TITLE, title="Minimal Plan")
        ])
        
        result_path = renderer.render(plan, temp_file)
        
        assert os.path.exists(result_path)
        assert result_path == temp_file
        
        # Verify it's a valid pptx
        prs = pptx.Presentation(result_path)
        assert len(prs.slides) == 1
        assert prs.slides[0].shapes.title.text == "Minimal Plan"

    def test_render_all_slide_types(self, renderer, temp_file):
        plan = ArtifactPlan(slides=[
            SlideModel(slide_type=SlideType.TITLE, title="Title"),
            SlideModel(slide_type=SlideType.CONTENT, title="Content", content=["1", "2"]),
            SlideModel(slide_type=SlideType.CONCEPT, title="Concept", content=["c1"]),
            SlideModel(slide_type=SlideType.EXAMPLE, title="Example", content=["e1"]),
            SlideModel(slide_type=SlideType.QUESTION, title="Question", content=["q1"]),
        ])
        
        renderer.render(plan, temp_file)
        
        prs = pptx.Presentation(temp_file)
        assert len(prs.slides) == 5

    def test_render_unicode_and_long_text(self, renderer, temp_file):
        long_text = "A" * 3000
        plan = ArtifactPlan(slides=[
            SlideModel(slide_type=SlideType.CONTENT, title="Unicode 🚀 β Σ ≈", content=["Test 123", long_text])
        ])
        
        renderer.render(plan, temp_file)
        
        prs = pptx.Presentation(temp_file)
        slide = prs.slides[0]
        assert slide.shapes.title.text == "Unicode 🚀 β Σ ≈"
        assert len(slide.placeholders) > 1
        # It shouldn't crash, and the text should be stored.

    def test_provenance_in_speaker_notes(self, renderer, temp_file):
        plan = ArtifactPlan(slides=[
            SlideModel(
                slide_type=SlideType.CONTENT, 
                title="Provenance", 
                content=["test"],
                speaker_notes="Here are some notes.",
                source_node_ids=["node1", "node2"],
                evidence_ids=["ev1"]
            )
        ])
        
        renderer.render(plan, temp_file)
        
        prs = pptx.Presentation(temp_file)
        slide = prs.slides[0]
        
        notes = slide.notes_slide.notes_text_frame.text
        assert "[Speaker Notes]" in notes
        assert "Here are some notes." in notes
        assert "[Provenance]" in notes
        assert "Sources: node1, node2" in notes
        assert "Evidence: ev1" in notes

    def test_rendering_failure_cleanup(self, renderer, temp_file):
        # We can force an exception during rendering by passing an invalid path to save
        plan = ArtifactPlan(slides=[
            SlideModel(slide_type=SlideType.TITLE, title="Fail")
        ])
        
        invalid_path = "/invalid/dir/that/doesnt/exist/file.pptx"
        tmp_path = f"{invalid_path}.tmp"
        
        with pytest.raises(Exception):
            renderer.render(plan, invalid_path)
            
        assert not os.path.exists(tmp_path)
        assert not os.path.exists(invalid_path)
